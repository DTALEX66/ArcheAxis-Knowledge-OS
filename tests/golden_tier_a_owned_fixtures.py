"""Generate project-authored Office fixture bytes for the G0 Tier-A corpus."""

from __future__ import annotations

from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile


def write_office_fixtures(destination: Path) -> None:
    """Write minimal no-personal-data DOCX, PPTX and XLSX evidence fixtures."""
    from pptx import Presentation
    from openpyxl import Workbook

    docx = destination / "golden-docx-anchor.docx"
    with ZipFile(docx, "w", ZIP_DEFLATED) as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "word/document.xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body>
<w:p><w:pPr><w:pStyle w:val="Heading1"/></w:pPr><w:r><w:t>Document evidence anchor</w:t></w:r></w:p>
<w:p><w:r><w:t>This project-authored fixture contains no personal data.</w:t></w:r></w:p>
<w:sectPr/></w:body></w:document>""",
        )

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Slide evidence anchor"
    slide.placeholders[1].text = "This project-authored fixture contains no personal data."
    presentation.save(destination / "golden-pptx-anchor.pptx")

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Evidence"
    sheet["A1"] = "Sheet evidence anchor"
    sheet["A2"] = "This project-authored fixture contains no personal data."
    workbook.save(destination / "golden-xlsx-anchor.xlsx")
    workbook.close()


if __name__ == "__main__":
    write_office_fixtures(Path(__file__).resolve().parent / "fixtures" / "golden")
