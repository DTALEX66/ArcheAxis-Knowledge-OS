"""AXW-023B: structured PPTX adapter (honest degradation + slide anchors).

Converts a .pptx into structured blocks with slide-level anchors using
python-pptx (already in the dependency tree via markitdown[pptx]). Slide
order is preserved; text, notes and tables are extracted per slide; media
references are reported as loss notes (visual/media content is not
extracted to text). Fails closed when python-pptx is unavailable.
"""

from __future__ import annotations

import hashlib
from pathlib import Path
from typing import Any

from shared.adapter_contract import AdapterResult


def _sha256(file_path: str | Path) -> str:
    h = hashlib.sha256()
    with open(file_path, "rb") as fh:
        for chunk in iter(lambda: fh.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def _extract_slides(path: Path) -> tuple[list[dict[str, Any]], list[str]]:
    """Return (blocks, loss_notes) with slide-ordered text/notes/tables.

    Media (pictures / videos / charts) cannot be represented as text; each
    occurrence is recorded as a loss note so the LossReport is honest about
    what the adapter dropped.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    blocks: list[dict[str, Any]] = []
    loss: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        media_kinds: list[str] = []

        for shape in slide.shapes:
            shape_type = getattr(shape, "shape_type", None)
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                media_kinds.append("picture")
                continue
            if shape.has_text_frame:
                text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                if text.strip():
                    slide_parts.append(text)
            if getattr(shape, "has_table", False) and shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    rows.append(" | ".join(cells))
                slide_parts.append("\n".join(rows))

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"[notes] {notes}")

        if slide_parts:
            blocks.append(
                {
                    "kind": "slide",
                    "text": "\n\n".join(slide_parts),
                    "anchor": {"slide_index": idx, "slide_number": idx},
                }
            )
        else:
            blocks.append(
                {"kind": "slide-empty", "text": "", "anchor": {"slide_index": idx, "slide_number": idx}}
            )

        if media_kinds:
            loss.append(f"slide {idx}: dropped media ({', '.join(sorted(set(media_kinds)))})")

    return blocks, loss


def convert_pptx(file_path: str | Path) -> AdapterResult:
    """Convert a .pptx to structured slide blocks.

    Returns markdown-like text in ``content`` and per-slide structure in
    ``metadata["blocks"]`` plus ``metadata["loss_notes"]``.
    """
    path = Path(file_path)
    if not path.is_file():
        return AdapterResult(success=False, content="", engine="pptx-adapter", error="file not found")

    try:
        from pptx import Presentation  # noqa: F401  (python-pptx present?)
    except ImportError:
        return AdapterResult(
            success=False,
            content="",
            engine="pptx-adapter",
            error="PPTX conversion requires python-pptx; install markitdown[pptx].",
        )

    try:
        blocks, loss = _extract_slides(path)
    except Exception as exc:  # pragma: no cover - host-tooling dependent
        return AdapterResult(
            success=False,
            content="",
            engine="pptx-adapter",
            error=f"python-pptx conversion failed: {exc}",
        )

    if not any(b["kind"] != "slide-empty" for b in blocks):
        return AdapterResult(
            success=False,
            content="",
            engine="pptx-adapter",
            error="PPTX conversion returned no content; treat as degraded.",
        )

    text = "\n\n".join(b["text"] for b in blocks if b["text"])
    return AdapterResult(
        success=True,
        content=text,
        engine="pptx-adapter",
        metadata={
            "char_count": len(text),
            "block_count": len(blocks),
            "blocks": blocks,
            "loss_notes": loss,
        },
    )


def convert_pptx_to_run(
    file_path: str | Path,
    db: str | Path,
    source_name: str | None = None,
    version: int = 1,
) -> dict[str, Any]:
    """Convert a .pptx and persist a ConversionRun; returns run metadata."""
    result = convert_pptx(file_path)
    if not result.success:
        raise RuntimeError(result.error or "PPTX conversion failed")

    from app.ingestion.conversion_run import create_conversion_run, store_conversion_run

    raw_sha = _sha256(file_path)
    blocks: list[dict[str, Any]] = result.metadata.get("blocks") or []
    loss = result.metadata.get("loss_notes") or []
    run = create_conversion_run(
        raw_sha256=raw_sha,
        source_name=source_name or Path(file_path).name,
        blocks=blocks,
        engine=result.engine,
        version=version,
    )
    store_conversion_run(db, run)
    return {
        "run_id": run.run_id,
        "document_id": run.document.document_id,
        "block_count": len(blocks),
        "loss_notes": loss,
    }
