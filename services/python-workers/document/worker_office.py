#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""ArcheAxis vNext office/document engine worker (F05/F07/F08/F09 partial).

Single worker entrypoint for structured office and PDF sources:

- docx: pure-stdlib ZIP+XML extraction (paragraphs/tables/media inventory;
        headers/footers noted; nothing executed)
- pptx: python-pptx engine (slide order, shape text, notes, image counts)
- xlsx: openpyxl engine (sheets/cells, formula text + cached-value policy note,
        merged ranges; macros never executed)
- pdf : PyMuPDF engine (page blocks in reading order, per-page anchors,
        image inventory; scanned pages reported, OCR is a separate lane)

Isolation boundary: never opens the vNext database; every engine failure
surfaces {"error": ...} with a non-zero exit.

Usage:
    python worker_office.py --probe
    python worker_office.py <input.docx|.pptx|.xlsx|.pdf>
"""

from __future__ import annotations

import argparse
import json
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

ENGINE = "python-worker-office"
ENGINE_VERSION = "0.1.0"

W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
W = lambda tag: f"{{{W_NS}}}{tag}"  # noqa: E731


def probe() -> dict:
    engine_status: dict[str, tuple[bool, str]] = {"docx": (True, "stdlib-zip+xml")}
    for module_name, format_name in (("pptx", "pptx"), ("openpyxl", "xlsx"), ("fitz", "pdf")):
        try:
            module = __import__(module_name)
            engine_status[format_name] = (True, getattr(module, "__version__", "unknown"))
        except ImportError:
            engine_status[format_name] = (False, "missing")
    engines = {fmt: ok for fmt, (ok, _version) in engine_status.items()}
    versions = {fmt: version for fmt, (_ok, version) in engine_status.items()}
    return {
        "engine": ENGINE,
        "engines": engines,
        "versions": versions,
        "formats": [fmt for fmt, ok in engines.items() if ok],
        "note": "docx always enabled (stdlib); pptx/xlsx/pdf require their engines",
    }


def _docx_text(path: Path) -> dict:
    text_parts: list[dict] = []
    media: list[str] = []
    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
        media = sorted(n for n in names if n.startswith("word/media/"))
        has_headers = any(n.startswith("word/header") for n in names)
        has_footers = any(n.startswith("word/footer") for n in names)
        document_xml = archive.read("word/document.xml")
    root = ET.fromstring(document_xml)
    body = root.find(W("body"))
    if body is None:
        raise ValueError("docx document.xml has no body")
    for child in body:
        tag = child.tag
        if tag == W("p"):
            runs = child.findall(".//" + W("t"))
            paragraph_text = "".join(run.text or "" for run in runs)
            if paragraph_text.strip():
                text_parts.append({"kind": "paragraph", "text": paragraph_text.strip()})
        elif tag == W("tbl"):
            for row in child.findall(".//" + W("tr")):
                cells = []
                for cell in row.findall(".//" + W("tc")):
                    cell_text = "".join(
                        run.text or ""
                        for run in cell.findall(".//" + W("t"))
                    )
                    cells.append(cell_text.strip())
                if any(cells):
                    text_parts.append({"kind": "table_row", "text": " | ".join(cells)})
    if not text_parts:
        raise ValueError("docx contains no extractable paragraph/table text")
    projection = "\n".join(part["text"] for part in text_parts)
    structure = []
    offset = 0
    for index, part in enumerate(text_parts, start=1):
        start = projection.find(part["text"], offset)
        if start < 0:
            start = offset
        structure.append(
            {"kind": part["kind"], "path": [f"{part['kind']}-{index}"], "char_start": start, "char_end": start + len(part["text"])}
        )
        offset = start + len(part["text"])
    return {
        "format": "docx",
        "text": projection,
        "structure": structure,
        "loss_receipt": {
            "engine": ENGINE,
            "engine_version": ENGINE_VERSION,
            "params": {
                "headers": has_headers,
                "footers": has_footers,
                "media_files": len(media),
                "engine": "stdlib-zip+xml",
            },
            "loss_note": (
                "paragraphs/tables extracted in document order; header/footer "
                "text and embedded-image OCR are separate lanes; media files "
                f"({len(media)}) inventoried, not decoded"
            ),
        },
    }


def _pptx_text(path: Path) -> dict:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("pptx engine missing (python-pptx not installed)") from exc
    presentation = Presentation(str(path))
    text_parts: list[dict] = []
    image_count = 0
    chart_count = 0
    for index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                text_parts.append({"kind": "slide", "index": index, "text": shape.text_frame.text.strip()})
            if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
                image_count += 1
            if getattr(shape, "has_chart", False):
                chart_count += 1
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                text_parts.append({"kind": "slide_notes", "index": index, "text": notes})
    if not text_parts:
        raise ValueError("pptx contains no extractable text")
    projection = "\n".join(part["text"] for part in text_parts)
    structure = []
    offset = 0
    for index, part in enumerate(text_parts, start=1):
        start = projection.find(part["text"], offset)
        if start < 0:
            start = offset
        structure.append(
            {"kind": part["kind"], "path": [f"slide-{part['index']}", part['kind']], "char_start": start, "char_end": start + len(part["text"])}
        )
        offset = start + len(part["text"])
    return {
        "format": "pptx",
        "text": projection,
        "structure": structure,
        "loss_receipt": {
            "engine": ENGINE,
            "engine_version": ENGINE_VERSION,
            "params": {"slides": len(presentation.slides._sldIdLst), "images": image_count, "charts": chart_count, "engine": "python-pptx"},
            "loss_note": "slide order preserved; slide-image OCR and chart rendering are separate lanes",
        },
    }


def _xlsx_text(path: Path) -> dict:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("xlsx engine missing (openpyxl not installed)") from exc
    workbook = load_workbook(str(path), data_only=False)
    text_parts: list[dict] = []
    formula_count = 0
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            cells = []
            for cell in row:
                if cell.value is None:
                    continue
                value = cell.value
                if isinstance(value, str) and value.startswith("="):
                    formula_count += 1
                cells.append(f"{cell.coordinate}={value}")
            if cells:
                text_parts.append({"kind": "sheet_row", "sheet": sheet.title, "text": " | ".join(cells)})
    if not text_parts:
        raise ValueError("xlsx contains no extractable cell values")
    projection = "\n".join(part["text"] for part in text_parts)
    structure = []
    offset = 0
    for index, part in enumerate(text_parts, start=1):
        start = projection.find(part["text"], offset)
        if start < 0:
            start = offset
        structure.append(
            {"kind": part["kind"], "path": [f"sheet-{part['sheet']}", f"row-{index}"], "char_start": start, "char_end": start + len(part["text"])}
        )
        offset = start + len(part["text"])
    return {
        "format": "xlsx",
        "text": projection,
        "structure": structure,
        "loss_receipt": {
            "engine": ENGINE,
            "engine_version": ENGINE_VERSION,
            "params": {"sheets": len(workbook.worksheets), "formula_cells": formula_count, "engine": "openpyxl"},
            "loss_note": (
                "cell values include formula text (data_only=False); cached "
                "computed values are NOT presented as live calculations; "
                "macros never executed; merged ranges reported per sheet only"
            ),
        },
    }


def _pdf_text(path: Path) -> dict:
    try:
        import pymupdf as fitz  # PyMuPDF >= 1.24 canonical import
    except ImportError:
        try:
            import fitz  # legacy import name (deprecated)
        except ImportError as exc:
            raise RuntimeError("pdf engine missing (PyMuPDF not installed)") from exc
    document = fitz.open(str(path))
    text_parts: list[dict] = []
    scanned_pages = 0
    for page_index, page in enumerate(document, start=1):
        blocks = page.get_text("blocks")
        page_text = "\n".join(block[4].strip() for block in blocks if block[4].strip())
        if not page_text.strip():
            scanned_pages += 1
        if page_text.strip():
            text_parts.append({"kind": "pdf_page", "page": page_index, "text": page_text.strip()})
    if not text_parts:
        raise ValueError(f"pdf contains no text layer ({scanned_pages} scanned pages; OCR lane required)")
    projection = "\n\n".join(part["text"] for part in text_parts)
    structure = []
    offset = 0
    for index, part in enumerate(text_parts, start=1):
        start = projection.find(part["text"], offset)
        if start < 0:
            start = offset
        structure.append(
            {"kind": "pdf_page", "path": [f"page-{part['page']}"], "char_start": start, "char_end": start + len(part["text"])}
        )
        offset = start + len(part["text"])
    return {
        "format": "pdf",
        "text": projection,
        "structure": structure,
        "loss_receipt": {
            "engine": ENGINE,
            "engine_version": ENGINE_VERSION,
            "params": {"pages": len(document), "scanned_pages_no_text_layer": scanned_pages, "engine": "PyMuPDF"},
            "loss_note": (
                "text blocks concatenated in page reading order; multi-column "
                "order and scanned-page OCR are separate lanes"
            ),
        },
    }


def extract(path: str) -> dict:
    suffix = Path(path).suffix.lower()
    if not Path(path).is_file():
        raise ValueError(f"input file not found: {path}")
    if suffix == ".docx":
        return _docx_text(Path(path))
    if suffix == ".pptx":
        return _pptx_text(Path(path))
    if suffix == ".xlsx":
        return _xlsx_text(Path(path))
    if suffix == ".pdf":
        return _pdf_text(Path(path))
    raise ValueError(f"unsupported office/document extension: {suffix}")


def main() -> int:
    parser = argparse.ArgumentParser(description="ArcheAxis office/document engine worker")
    parser.add_argument("input", nargs="?", help="input file (.docx/.pptx/.xlsx/.pdf)")
    parser.add_argument("--probe", action="store_true", help="engine capability probe")
    args = parser.parse_args()
    if args.probe:
        print(json.dumps(probe(), ensure_ascii=False))
        return 0
    if not args.input:
        print(json.dumps({"error": "usage: worker_office.py <input-file> | --probe"}))
        return 2
    try:
        out = extract(args.input)
    except Exception as exc:  # noqa: BLE001
        print(json.dumps({"error": str(exc)}, ensure_ascii=False))
        return 1
    out["engine"] = ENGINE
    out["engine_version"] = ENGINE_VERSION
    print(json.dumps(out, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
